#!/usr/bin/env python3
"""Build an editable layered PPTX from an Illustrator artboard export manifest."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

EMU_PER_PT = 12700

# Keep defaults conservative. Project-specific aliases should be supplied with
# --font-map-json rather than committed to a public repository.
FONT_NAME_ALIASES = {
    "ArialMT": "Arial",
    "Arial-BoldMT": "Arial",
    "Arial-ItalicMT": "Arial",
    "TimesNewRomanPSMT": "Times New Roman",
    "TimesNewRomanPS-BoldMT": "Times New Roman",
    "TimesNewRomanPS-ItalicMT": "Times New Roman",
    "SimSun": "SimSun",
    "MicrosoftYaHei": "Microsoft YaHei",
}

FONT_FALLBACK = "Arial"


def pt_to_inches(value: float) -> float:
    return float(value) / 72.0


def pt_to_emu(value: float) -> int:
    return int(round(float(value) * EMU_PER_PT))


def clean_font_value(value: str | None) -> str:
    return str(value or "").strip()


def rgb_from_hex(value: str | None) -> RGBColor:
    value = (value or "FFFFFF").strip().lstrip("#")
    if len(value) != 6:
        value = "FFFFFF"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def clean_hex(value: str | None) -> str:
    value = (value or "FFFFFF").strip().lstrip("#").upper()
    if len(value) != 6 or any(ch not in "0123456789ABCDEF" for ch in value):
        return "FFFFFF"
    return value


def remove_text_fill(rpr) -> None:
    fill_tags = {qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill"), qn("a:pattFill"), qn("a:grpFill")}
    for child in list(rpr):
        if child.tag in fill_tags:
            rpr.remove(child)


def set_run_gradient_fill(run, gradient: dict[str, Any] | None) -> None:
    if not gradient:
        return
    stops = gradient.get("stops") or []
    if len(stops) < 2:
        return

    rpr = run._r.get_or_add_rPr()
    remove_text_fill(rpr)

    grad_fill = OxmlElement("a:gradFill")
    grad_fill.set("rotWithShape", "1")
    gs_lst = OxmlElement("a:gsLst")
    for stop in stops:
        gs = OxmlElement("a:gs")
        position = int(float(stop.get("position", 0)))
        gs.set("pos", str(max(0, min(100000, position))))
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", clean_hex(stop.get("color")))
        gs.append(srgb)
        gs_lst.append(gs)
    grad_fill.append(gs_lst)

    lin = OxmlElement("a:lin")
    angle = float(gradient.get("angle") or 0)
    lin.set("ang", str(int(round(angle * 60000))))
    lin.set("scaled", "1")
    grad_fill.append(lin)
    rpr.append(grad_fill)


def mapped_font(
    font_name: str | None,
    font_family: str | None = None,
    font_full_name: str | None = None,
    font_style: str | None = None,
) -> str:
    candidates = [
        clean_font_value(font_full_name),
        clean_font_value(font_family),
        clean_font_value(font_name),
    ]
    for candidate in candidates:
        if candidate and candidate in FONT_NAME_ALIASES:
            return FONT_NAME_ALIASES[candidate]

    # Prefer Illustrator's family name because PPT usually expects a family face,
    # while fontName is often a PostScript identifier.
    for candidate in (candidates[1], candidates[0], candidates[2]):
        if candidate:
            return candidate
    return FONT_FALLBACK


def apply_font_style(run, font_name: str | None, font_style: str | None) -> None:
    combined = f"{font_name or ''} {font_style or ''}".lower()
    run.font.bold = any(token in combined for token in ("bold", "semibold", "demibold", "heavy", "black"))
    run.font.italic = any(token in combined for token in ("italic", "oblique"))


def set_east_asian_font(run, font_name: str) -> None:
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = rpr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            rpr.append(element)
        element.set("typeface", font_name)


def clear_shape_style(shape) -> None:
    shape.fill.background()
    shape.line.fill.background()
    sp_pr = shape._element.spPr
    for child in list(sp_pr):
        if child.tag.endswith("glow") or child.tag.endswith("blur"):
            sp_pr.remove(child)


def alignment_from_illustrator(value: str | None):
    text = (value or "").upper()
    if "CENTER" in text:
        return PP_ALIGN.CENTER
    if "RIGHT" in text:
        return PP_ALIGN.RIGHT
    if "JUSTIFY" in text:
        return getattr(PP_ALIGN, "JUSTIFY", PP_ALIGN.LEFT)
    return PP_ALIGN.LEFT


def set_paragraph_defaults(paragraph, alignment) -> None:
    paragraph.alignment = alignment
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0


def iter_text_chunks(runs: Iterable[dict[str, Any]]) -> Iterable[tuple[str, dict[str, Any], bool]]:
    first = True
    for run in runs:
        text = str(run.get("text") or "")
        parts = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
        for index, part in enumerate(parts):
            new_para = not first and index > 0
            yield part, run, new_para
            first = False


def add_text(shape, item: dict[str, Any]) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)

    body_pr = text_frame._txBody.bodyPr
    body_pr.set("wrap", "square")
    body_pr.set("anchor", "t")

    alignment = alignment_from_illustrator(item.get("justification"))
    paragraph = text_frame.paragraphs[0]
    set_paragraph_defaults(paragraph, alignment)

    runs = item.get("runs") or [item]
    wrote_any = False
    for text, run_data, new_para in iter_text_chunks(runs):
        if new_para:
            paragraph = text_frame.add_paragraph()
            set_paragraph_defaults(paragraph, alignment)
        if text == "" and wrote_any:
            continue
        ppt_run = paragraph.add_run()
        ppt_run.text = text
        font_name = mapped_font(
            run_data.get("fontName"),
            run_data.get("fontFamily"),
            run_data.get("fontFullName"),
            run_data.get("fontStyle"),
        )
        ppt_run.font.name = font_name
        ppt_run.font.size = Pt(float(run_data.get("size") or item.get("size") or 24))
        ppt_run.font.color.rgb = rgb_from_hex(run_data.get("fillColor") or item.get("fillColor"))
        apply_font_style(ppt_run, run_data.get("fontName"), run_data.get("fontStyle"))
        set_east_asian_font(ppt_run, font_name)
        set_run_gradient_fill(ppt_run, item.get("textGradient"))
        wrote_any = True

    if not wrote_any:
        ppt_run = paragraph.add_run()
        ppt_run.text = str(item.get("contents") or "")
        font_name = mapped_font(item.get("fontName"), item.get("fontFamily"), item.get("fontFullName"), item.get("fontStyle"))
        ppt_run.font.name = font_name
        ppt_run.font.size = Pt(float(item.get("size") or 24))
        ppt_run.font.color.rgb = rgb_from_hex(item.get("fillColor"))
        apply_font_style(ppt_run, item.get("fontName"), item.get("fontStyle"))
        set_east_asian_font(ppt_run, font_name)
        set_run_gradient_fill(ppt_run, item.get("textGradient"))


def merged_text_bounds(item: dict[str, Any]) -> dict[str, float]:
    geometric = item.get("geometricBounds") or item.get("bounds") or item.get("visibleBounds") or {}
    visible = item.get("visibleBounds") or geometric
    left = min(float(geometric.get("x", 0)), float(visible.get("x", 0)))
    top = min(float(geometric.get("y", 0)), float(visible.get("y", 0)))
    right = max(
        float(geometric.get("x", 0)) + float(geometric.get("w", 0)),
        float(visible.get("x", 0)) + float(visible.get("w", 0)),
    )
    bottom = max(
        float(geometric.get("y", 0)) + float(geometric.get("h", 0)),
        float(visible.get("y", 0)) + float(visible.get("h", 0)),
    )
    size = float(item.get("size") or 24)
    pad_x = max(6.0, size * 0.45)
    pad_y = max(3.0, size * 0.15)
    return {
        "x": max(0.0, left - pad_x / 2),
        "y": max(0.0, top - pad_y / 2),
        "w": max(1.0, right - left + pad_x),
        "h": max(1.0, bottom - top + pad_y),
    }


def load_font_aliases(path: Path | None) -> None:
    if not path:
        return
    data = json.loads(path.read_text(encoding="utf-8-sig"))
    if not isinstance(data, dict):
        raise ValueError("font alias JSON must be an object: {sourceName: pptFontName}")
    FONT_NAME_ALIASES.update({str(k): str(v) for k, v in data.items()})


def main() -> int:
    global FONT_FALLBACK
    parser = argparse.ArgumentParser(description="Build editable PPTX from Illustrator manifest.json")
    parser.add_argument("--manifest", required=True, type=Path, help="Path to manifest.json exported by export_artboard.jsx")
    parser.add_argument("--out", required=True, type=Path, help="Output .pptx path")
    parser.add_argument("--keep-image-order", action="store_true", help="Do not reverse image layer order when adding to PPT")
    parser.add_argument("--font-map-json", type=Path, help="Optional JSON aliases for source font names to PPT font names")
    parser.add_argument("--font-fallback", default=FONT_FALLBACK, help="Fallback font only when AI provides no font name")
    args = parser.parse_args()

    FONT_FALLBACK = args.font_fallback
    load_font_aliases(args.font_map_json)

    manifest_path = args.manifest.resolve()
    export_dir = manifest_path.parent
    manifest = json.loads(manifest_path.read_text(encoding="utf-8-sig"))

    artboard = manifest["artboardRect"]
    width_pt = float(artboard[2]) - float(artboard[0])
    height_pt = float(artboard[1]) - float(artboard[3])

    prs = Presentation()
    prs.slide_width = Inches(pt_to_inches(width_pt))
    prs.slide_height = Inches(pt_to_inches(height_pt))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    images = manifest.get("images", [])
    image_iterable = images if args.keep_image_order else list(reversed(images))
    for image in image_iterable:
        bounds = image["bounds"]
        image_path = export_dir / image["file"]
        if not image_path.exists():
            raise FileNotFoundError(image_path)
        slide.shapes.add_picture(
            str(image_path),
            pt_to_emu(bounds["x"]),
            pt_to_emu(bounds["y"]),
            width=pt_to_emu(bounds["w"]),
            height=pt_to_emu(bounds["h"]),
        )

    for item in manifest.get("texts", []):
        bounds = merged_text_bounds(item)
        shape = slide.shapes.add_textbox(
            pt_to_emu(bounds["x"]),
            pt_to_emu(bounds["y"]),
            pt_to_emu(bounds["w"]),
            pt_to_emu(bounds["h"]),
        )
        clear_shape_style(shape)
        add_text(shape, item)

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"Wrote {args.out}")
    print(f"Slide size: {width_pt / 72 * 2.54:.2f} x {height_pt / 72 * 2.54:.2f} cm")
    print(f"Pictures: {len(images)}, text boxes: {len(manifest.get('texts', []))}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
