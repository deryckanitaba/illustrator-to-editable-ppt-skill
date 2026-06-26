#!/usr/bin/env python3
"""Verify a generated editable Illustrator-to-PPTX slide."""

from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

EMU_PER_INCH = 914400


def cm_from_emu(value: int) -> float:
    return value / EMU_PER_INCH * 2.54


def read_xml(pptx: Path) -> str:
    chunks: list[str] = []
    with zipfile.ZipFile(pptx) as archive:
        for name in archive.namelist():
            if name.endswith(".xml"):
                chunks.append(archive.read(name).decode("utf-8", errors="ignore"))
    return "\n".join(chunks)


def expected_size_from_manifest(manifest: Path) -> tuple[float, float] | None:
    if not manifest:
        return None
    data = json.loads(manifest.read_text(encoding="utf-8-sig"))
    ab = data["artboardRect"]
    width_cm = (float(ab[2]) - float(ab[0])) / 72 * 2.54
    height_cm = (float(ab[1]) - float(ab[3])) / 72 * 2.54
    return width_cm, height_cm


def summarize_fonts(xml: str) -> Counter[str]:
    fonts = re.findall(r'typeface="([^"]+)"', xml)
    return Counter(fonts)


def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect a generated editable PPTX")
    parser.add_argument("pptx", type=Path, help="PPTX file to verify")
    parser.add_argument("--manifest", type=Path, help="Optional Illustrator export manifest.json")
    parser.add_argument("--expect-pictures", type=int, help="Expected picture count on the first slide")
    parser.add_argument("--expect-texts", type=int, help="Expected text box count on the first slide")
    parser.add_argument("--expect-font", help="Optional font marker expected in PPT XML")
    parser.add_argument("--show-fonts", action="store_true", help="Print top font typefaces found in PPT XML")
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    if not prs.slides:
        print("ERROR: no slides found")
        return 1
    slide = prs.slides[0]
    pictures = sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    texts = sum(1 for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    width_cm = cm_from_emu(prs.slide_width)
    height_cm = cm_from_emu(prs.slide_height)

    xml = read_xml(args.pptx)
    metrics = {
        "glow": xml.count("<a:glow"),
        "blur": xml.count("<a:blur"),
        "wrap_square": xml.count('wrap="square"'),
        "wrap_none": xml.count('wrap="none"'),
        "center_alignment": xml.count('algn="ctr"'),
        "right_alignment": xml.count('algn="r"'),
    }
    if args.expect_font:
        metrics["expected_font_marker"] = xml.count(args.expect_font)

    print(f"PPTX: {args.pptx}")
    print(f"Slide size: {width_cm:.2f} x {height_cm:.2f} cm")
    print(f"Pictures: {pictures}")
    print(f"Editable text boxes: {texts}")
    for key, value in metrics.items():
        print(f"{key}: {value}")
    if args.show_fonts:
        for font, count in summarize_fonts(xml).most_common(20):
            print(f"font: {font} ({count})")

    failures: list[str] = []
    if args.manifest:
        expected_size = expected_size_from_manifest(args.manifest)
        if expected_size:
            ew, eh = expected_size
            if abs(width_cm - ew) > 0.02 or abs(height_cm - eh) > 0.02:
                failures.append(f"slide size mismatch, expected {ew:.2f} x {eh:.2f} cm")
        data = json.loads(args.manifest.read_text(encoding="utf-8-sig"))
        args.expect_pictures = args.expect_pictures if args.expect_pictures is not None else len(data.get("images", []))
        args.expect_texts = args.expect_texts if args.expect_texts is not None else len(data.get("texts", []))

    if args.expect_pictures is not None and pictures != args.expect_pictures:
        failures.append(f"picture count mismatch, expected {args.expect_pictures}")
    if args.expect_texts is not None and texts != args.expect_texts:
        failures.append(f"text count mismatch, expected {args.expect_texts}")
    if metrics["glow"] or metrics["blur"]:
        failures.append("glow/blur effect markers found")
    if metrics["wrap_none"]:
        failures.append('wrap="none" markers found')
    if texts and metrics["wrap_square"] < texts:
        failures.append("not every text box has square text wrapping")
    if args.expect_font and metrics.get("expected_font_marker", 0) == 0:
        failures.append(f"font marker not found: {args.expect_font}")

    if failures:
        print("FAIL:")
        for failure in failures:
            print(f"- {failure}")
        return 1
    print("PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
